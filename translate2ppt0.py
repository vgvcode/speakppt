from pptx import Presentation

def createTitleSlide(p, layout, content):
    slide = p.slides.add_slide(layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = content['title']
    subtitle.text = content['subtitle']

def createBulletSlide(p, layout, content):
    slide = p.slides.add_slide(layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = content['title']

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    id = 0
    for b in content['bullets']:
        id = id + 1
        if id == 1:
            tf.text = b['l1']
        else:
            p = tf.add_paragraph()
            p.text = b['l1']

        if 'l2' in b:
            p = tf.add_paragraph()
            p.text = b['l2']
            p.level = 1
        if 'l3' in b:
            p = tf.add_paragraph()
            p.text = b['l3']
            p.level = 2


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_slide_content = {
    "title": "Hello World!",
    "subtitle": "By VGV"
}

bullet_slide_layout = prs.slide_layouts[1]
bullet_slide_content = {
    "title": 'Adding a bullet slide',
    "bullets": [
        {
            "l1": "First bullet level 1",
            "l2": "First bullet level 2",
            "l3": "First bullet level 3"
        },
        {
            "l1": "Second bullet level 1",
            "l2": "Second bullet level 2",
            "l3": "Second bullet level 3"
        }
    ]
}

createTitleSlide(prs, title_slide_layout, title_slide_content)
createBulletSlide(prs, bullet_slide_layout, bullet_slide_content)

fileName = "test.pptx"
prs.save(fileName)
