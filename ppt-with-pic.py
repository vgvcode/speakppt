import sys
from pptx import Presentation
from pptx.util import Inches

img_path = 'aditi-1.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

#left = top = Inches(1)
#pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(7.5)
top = Inches(4)
height = Inches(3)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('testpic.pptx')
