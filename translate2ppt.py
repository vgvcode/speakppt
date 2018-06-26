import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import glob
import random
import boto3
from urllib import parse
import json

s3 = boto3.resource('s3')

gPrs = None
gTitleSlideLayout = None
gBulletSlideLayout = None
gForwardSubstitutions = {"##" : u'\x81', "@@" : u'\x82', "!!" : u'\x83'}
gReverseSubstitutions = {u'\x81': "#", u'\x82': "@", u'\x83': '!'}
#default voice
gDefaultVoice = "Joanna"
gCurrentVoice = gDefaultVoice
#this is a collection of voices for this presentation
gVoices = [gDefaultVoice]

def doTitle(args):
    titleText = ' '.join(args).split('@')
    slide = getCurrentSlide()
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titleText[0]
    subtitle.text = titleText[1] if len(titleText) > 1 else ""

def doHeading(args):
    headingText = ' '.join(args)
    slide = getCurrentSlide()
    shapes = slide.shapes
    titleShape = shapes.title
    titleShape.text = headingText

def doVoices(args):
    global gVoices, gCurrentVoice
    gVoices = args
    gCurrentVoice = args[0]

def doParagraph(args):
    global gPrs
    para = ' '.join(args).split('@')
    para = doReverseSubstitutions(para)
    slide = getCurrentSlide()
    shapes = slide.shapes
    bodyShape = shapes.placeholders[1]
    tf = bodyShape.text_frame

    id = 0
    for b in para:
        id = id + 1
        if id == 1:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(24)

def doQuiz(args):
    doHeading(["A short quiz!"])
    doParagraph(args)

def doQuestions(args):
    doHeading(["Any Questions?"])
    doParagraph(args)

def doSummary(args):
    doHeading(["Summary"])
    doParagraph(args)

def doConclusion(args):
    doHeading(["Thank you!"])
    doParagraph(args)

def getNextVoice():
    try:
        if gCurrentVoice in gVoices:
            ix = gVoices.index(gCurrentVoice)
            ix = (ix + 1) % len(gVoices)
            nextVoice = gVoices[ix]
        else:
            nextVoice = gDefaultVoice
    except Exception:
        #no voice found in the list, setting to default voice
        print("Unexpected error @getNextVoice:", sys.exc_info()[0])
        nextVoice = gDefaultVoice
    finally:
        return nextVoice

def doTransition(args):
    global gCurrentVoice, gTitleSlideLayout, gBulletSlideLayout
    if (len(gPrs.slides) == 0):
        #first slide is a title slide
        slide = gPrs.slides.add_slide(gTitleSlideLayout)
    else:
        slide = gPrs.slides.add_slide(gBulletSlideLayout)
        gCurrentVoice = getNextVoice()

    #if a voice is specified in the transition command, honor it
    if len(args) > 0:
        gCurrentVoice = args[0]

    left = Inches(11)
    top = Inches(4)
    height = Inches(3)
    #make first char case insensitive
    path = "[" + gCurrentVoice[0].lower() + "|" + gCurrentVoice[0].upper() + "]" + gCurrentVoice[1:] + "*"
    imageList = glob.glob(path)
    ix = random.randrange(len(imageList))
    img_path = imageList[ix]
    print("Adding image to ppt: " + img_path)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    # move picture to background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)  # use the number that does the appropriate job

def doForwardSubstitutions(s):
    global gForwardSubstitutions
    for key in gForwardSubstitutions:
        s = s.replace(key, gForwardSubstitutions[key])
    return s

def doReverseSubstitutions(args):
    global gReverseSubstitutions
    newArgs = []
    for a in args:
        for key in gReverseSubstitutions:
            a = a.replace(key, gReverseSubstitutions[key])
        newArgs.append(a)
    return newArgs

def getCurrentSlide():
    global gPrs
    return gPrs.slides[len(gPrs.slides) - 1]

def processTag(section):
    words = section.split(' ')
    tag = words[0]
    args = words[1:]

    switcher = {
        'voices': lambda: doVoices(args),
        'title': lambda: doTitle(args),
        'paragraph': lambda: doParagraph(args),
        'heading': lambda: doHeading(args),
        'bullets': lambda: doParagraph(args),
        'transition': lambda: doTransition(args),
        'quiz': lambda: doQuiz(args),
        'questions': lambda: doQuestions(args),
        'summary': lambda: doSummary(args),
        'conclusion': lambda: doConclusion(args)
    }

    func = switcher.get(tag, lambda: '')
    return func()

def doTranslate(fileName):
    global gPrs

    f=open(fileName, "r")
    lines=f.readlines()
    txt=''.join(lines).replace('\n', ' ')
    txt=doForwardSubstitutions(txt)
    sections=txt.split("#")
    #print(sections)
    for s in sections:
        s = s.strip()
        if len(s) > 0:
            #print('Processing:' + s)
            processTag(s)

    outputFileName = fileName + ".pptm"
    gPrs.save(outputFileName)

# read file from s3
def readFromS3Bucket(b, s3Inf):
    #return bytes
    return s3.Object(b, s3Inf).get()['Body'].read()

def doTranslateInLambda(b, s3Inf):
    #b = vgvcode.speakppt
    #s3Inf = input/business.txt
    global gPrs

    #read the data file
    data = readFromS3Bucket(b, s3Inf).decode('utf-8')
    print("Read " + s3Inf + " from s3")

    lines = data.split('\n')
    txt=''.join(lines).replace('\n', ' ')
    txt=doForwardSubstitutions(txt)
    sections=txt.split("#")
    #print(sections)
    for s in sections:
        s = s.strip()
        if len(s) > 0:
            #print('Processing:' + s)
            processTag(s)

    #save the presentation to /tmp
    tmpOutf = s3Inf.replace("input", "/tmp", 1) + ".pptm"
    #tmpOutf = /tmp/business.txt.pptm
    gPrs.save(tmpOutf)
    print("Saved presentation to " + tmpOutf)

    #copy the presentation from /tmp to s3
    s3Outf =  s3Inf.replace("input", "output", 1) + ".pptm"
    #s3Outf = output/business.txt.pptm
    #write bytes
    s3.Object(b, s3Outf).put(Body=open(tmpOutf, 'rb'))
    print("Copied presentation to " + s3Outf)

def lambda_handler(event, context):
    global gPrs, gTitleSlideLayout, gBulletSlideLayout, gTemplateFile

    #bucket = vgvcode.speakppt
    #dataFile = input/business.txt
    #templateFile = input/pptm-empty-file-with-macros.pptm
    print(json.dumps(event))
    snsMsgObject = json.loads(event['Records'][0]['Sns']['Message'])
    bucket = snsMsgObject['Records'][0]['s3']['bucket']['name']
    key = snsMsgObject['Records'][0]['s3']['object']['key']

    #Object key may have spaces or unicode non-ASCII characters.
    dataFile = parse.unquote(key.replace('/\+/g', " "))
    templateFile = os.getenv("TEMPLATE_FILE", "input/pptm-empty-file-with-macros.pptm")

    print("Bucket:" + bucket + ", dataFile:" + dataFile + ", templateFile:" + templateFile)

    obj = s3.Object(bucket, templateFile)
    template = obj.get()['Body'].read()
    print("Read: " + templateFile + " from bucket: " + bucket)

    #save the template file in /tmp
    outputTemplatePath = templateFile.replace('input', '/tmp', 1)
    f = open(outputTemplatePath, "wb")
    f.write(template)
    f.close()
    print("Saved: " + outputTemplatePath)

    #create a presentation from the template
    gPrs = Presentation(outputTemplatePath)
    gTitleSlideLayout = gPrs.slide_layouts[0]
    gBulletSlideLayout = gPrs.slide_layouts[1]
    print("Started work on a ppt using: " + outputTemplatePath)

    doTranslateInLambda(bucket, dataFile)

def main():
    global gPrs, gTitleSlideLayout, gBulletSlideLayout
    if len(sys.argv) != 2:
        print("Usage:" + sys.argv[0] + " /path/to/manifest")
        exit()
    fileName = sys.argv[1]

    template = "./pptm-empty-file-with-macros.pptm"

    gPrs = Presentation(template)
    gTitleSlideLayout = gPrs.slide_layouts[0]
    gBulletSlideLayout = gPrs.slide_layouts[1]

    doTranslate(fileName)

if __name__ == "__main__":
    # execute only if run as a script
    main()
